#!/usr/bin/env python3
"""
Scalable Brain Project Presentation Generator
Creates a professional 5-page PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation with widescreen layout
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme - Professional dark theme
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)  # Dark navy
ACCENT_BLUE = RGBColor(0x00, 0x8f, 0xff)  # Bright blue
ACCENT_GREEN = RGBColor(0x00, 0xc8, 0x5c)  # Success green
ACCENT_RED = RGBColor(0xff, 0x4d, 0x4d)  # Alert red
LIGHT_TEXT = RGBColor(0xf0, 0xf0, 0xf0)  # Off-white
SUBTITLE_TEXT = RGBColor(0xb0, 0xb0, 0xc0)  # Light gray

def add_title_slide(prs):
    """Slide 1: Title and Overview"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    # Accent line at top
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_BLUE
    accent.line.fill.background()
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SCALABLE BRAIN"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(12.333), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Institutional-Grade Quantitative Trading System"
    p.font.size = Pt(32)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.6))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.text = "8-Layer AI-Powered Forex Trading Platform with Real-Time Market Intelligence"
    p.font.size = Pt(18)
    p.font.color.rgb = SUBTITLE_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # Key highlights
    highlights = [
        ("Real-Time Execution", ACCENT_GREEN),
        ("ML-Powered Filtering", ACCENT_BLUE),
        ("Multi-Layer Architecture", ACCENT_BLUE),
        ("OANDA Integration", ACCENT_BLUE)
    ]
    
    y_pos = 5.5
    for i, (text, color) in enumerate(highlights):
        x_pos = 1.0 + i * 3.0
        
        # Bullet
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos), Inches(y_pos), Inches(0.12), Inches(0.12))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = color
        bullet.line.fill.background()
        
        # Text
        text_box = slide.shapes.add_textbox(Inches(x_pos + 0.25), Inches(y_pos - 0.05), Inches(2.5), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_TEXT
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "April 2026"
    p.font.size = Pt(12)
    p.font.color.rgb = SUBTITLE_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_architecture_slide(prs):
    """Slide 2: 8-Layer System Architecture"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture: 8-Layer Runtime Model"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.LEFT
    
    # Architecture layers - left side
    layers = [
        ("Layer 0", "Strategy Qualification", "Backtest & qualify strategy variants", "0x4a90d9"),
        ("Layer 1", "Regime Detection", "Market state labels & context features", "0x5ba3e8"),
        ("Layer 2", "Signal Generation", "Rule-based indicator engine", "0x6cb6f7"),
        ("Layer 3", "ML Gatekeeper", "Train champion model artifacts", "0x7dc9ff"),
        ("Layer 4", "Live Execution", "Orchestration & risk gating", "0x00c85c"),
        ("Layer 5", "Telemetry API", "FastAPI dashboard & observability", "0xf5a623"),
        ("Layer 6", "Trade Auditor", "Outcome reconciliation", "0xe8913c"),
        ("Layer 7", "Broker Executor", "OANDA order placement", "0xff6b6b"),
    ]
    
    y_start = 1.3
    for i, (layer, name, desc, color_hex) in enumerate(layers):
        y_pos = y_start + i * 0.72
        color = RGBColor(int(color_hex[2:4], 16), int(color_hex[4:6], 16), int(color_hex[6:8], 16))
        
        # Layer box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(5.5), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Layer text
        text_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.08), Inches(5.1), Inches(0.45))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{layer}: {name}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = LIGHT_TEXT
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(6.2), Inches(y_pos + 0.12), Inches(6.5), Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = SUBTITLE_TEXT
    
    # Key principles box
    principles_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.4), Inches(12.333), Inches(0.9))
    principles_box.fill.solid()
    principles_box.fill.fore_color.rgb = RGBColor(0x2a, 0x2a, 0x4e)
    principles_box.line.color.rgb = ACCENT_BLUE
    
    principles_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12), Inches(0.7))
    tf = principles_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Core Principles: Explicit layer contracts • No downstream recomputation • Granularity preservation • Deterministic execution"
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_dataflow_slide(prs):
    """Slide 3: Data Flow and Core Components"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Core Components & Data Flow"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.LEFT
    
    # Data flow diagram - simplified representation
    components = [
        ("OANDA v20 API", "Real-time price streaming", 0.8, 0x5ba3e8),
        ("SQL Server Database", "Fact_Market_Prices, Fact_Signals, Fact_Live_Trades", 2.0, 0x6cb6f7),
        ("Signal Engine", "Rule-based indicator evaluation", 3.2, 0x7dc9ff),
        ("ML Gatekeeper", "RandomForest model filtering", 4.4, 0x00c85c),
        ("Risk Manager", "ATR-based position sizing, correlation gates", 5.6, 0xf5a623),
        ("Execution Engine", "Live trade placement via Layer 7", 6.8, 0xff6b6b),
    ]
    
    for name, desc, y_pos, color_hex in components:
        color_hex_str = hex(color_hex)[2:].zfill(6)
        color = RGBColor(int(color_hex_str[0:2], 16), int(color_hex_str[2:4], 16), int(color_hex_str[4:6], 16))
        
        # Component box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(y_pos), Inches(4.5), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Name
        name_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.1), Inches(4.1), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = LIGHT_TEXT
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.45), Inches(4.1), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0xe0, 0xe0, 0xe0)
        
        # Arrow down (except last)
        if y_pos < 6.5:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.0), Inches(y_pos + 0.95), Inches(0.5), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_BLUE
            arrow.line.fill.background()
    
    # Right side - Key features
    features_title = slide.shapes.add_textbox(Inches(6.5), Inches(1.0), Inches(6.333), Inches(0.5))
    tf = features_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Capabilities"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    features = [
        "• Multi-timeframe analysis (H1, H4)",
        "• Market regime detection (ATR/ADX)",
        "• 6 strategy variants (trend + range)",
        "• ML-based trade filtering",
        "• Correlation-aware risk management",
        "• Real-time WebSocket streaming",
        "• Automated trade auditing",
    ]
    
    y_pos = 1.6
    for feature in features:
        feat_box = slide.shapes.add_textbox(Inches(6.5), Inches(y_pos), Inches(6.333), Inches(0.4))
        tf = feat_box.text_frame
        p = tf.paragraphs[0]
        p.text = feature
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_TEXT
        y_pos += 0.45
    
    # Technology stack
    stack_title = slide.shapes.add_textbox(Inches(6.5), Inches(5.0), Inches(6.333), Inches(0.5))
    tf = stack_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Technology Stack"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    stack_items = [
        "Backend: Python 3.12, FastAPI, SQL Server",
        "Frontend: React 18, TypeScript, HTML5 Canvas",
        "ML: scikit-learn, pandas, numpy",
        "Broker: OANDA v20 API",
        "Data: WebSocket streaming, Redis caching"
    ]
    
    y_pos = 5.5
    for item in stack_items:
        item_box = slide.shapes.add_textbox(Inches(6.5), Inches(y_pos), Inches(6.333), Inches(0.35))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = SUBTITLE_TEXT
        y_pos += 0.38
    
    return slide

def add_charts_slide(prs):
    """Slide 4: Advanced Chart System"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Advanced Chart System: Custom HTML5 Canvas Solution"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.LEFT
    
    # Left column - Chart features
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.5))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Chart Features"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    features = [
        ("Real-Time Data", "WebSocket streaming from OANDA v20 API"),
        ("Technical Indicators", "SMA, EMA, RSI, MACD, Bollinger Bands, ADX, ATR"),
        ("Trade Visualization", "Entry markers, SL/TP lines, win/loss indicators"),
        ("Analysis Tools", "Trend lines, Fibonacci retracement, S/R levels"),
        ("Multi-Timeframe", "1m, 5m, 15m, 30m, 1h, 4h, 1d support"),
        ("Asset Filtering", "EUR_USD, GBP_USD, USD_JPY, AUD_USD, USD_CAD"),
    ]
    
    y_pos = 1.8
    for title, desc in features:
        # Title
        feat_title = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(5.8), Inches(0.35))
        tf = feat_title.text_frame
        p = tf.paragraphs[0]
        p.text = f"▸ {title}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = LIGHT_TEXT
        
        # Description
        feat_desc = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.35), Inches(5.5), Inches(0.3))
        tf = feat_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = SUBTITLE_TEXT
        
        y_pos += 0.75
    
    # Right column - Performance metrics
    right_title = slide.shapes.add_textbox(Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.5))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Performance Targets"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    metrics = [
        ("Initial Load", "< 2 seconds"),
        ("Real-time Latency", "< 50ms"),
        ("Frame Rate", "60 FPS"),
        ("Candle Capacity", "10,000+ without lag"),
        ("Memory Footprint", "< 50MB"),
        ("Indicator Calculation", "< 500ms"),
    ]
    
    y_pos = 1.8
    for metric, value in metrics:
        # Metric box
        metric_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(y_pos), Inches(5.8), Inches(0.55))
        metric_box.fill.solid()
        metric_box.fill.fore_color.rgb = RGBColor(0x2a, 0x2a, 0x4e)
        metric_box.line.color.rgb = ACCENT_BLUE
        
        # Metric name
        name_box = slide.shapes.add_textbox(Inches(7.2), Inches(y_pos + 0.08), Inches(3.0), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_TEXT
        
        # Value
        value_box = slide.shapes.add_textbox(Inches(10.0), Inches(y_pos + 0.08), Inches(2.6), Inches(0.4))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        p.alignment = PP_ALIGN.RIGHT
        
        y_pos += 0.7
    
    # Architecture note
    arch_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(5.8), Inches(5.8), Inches(1.3))
    arch_box.fill.solid()
    arch_box.fill.fore_color.rgb = RGBColor(0x2a, 0x2a, 0x4e)
    arch_box.line.color.rgb = ACCENT_BLUE
    
    arch_title = slide.shapes.add_textbox(Inches(7.2), Inches(5.9), Inches(5.4), Inches(0.4))
    tf = arch_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Architecture Highlights"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    arch_text = slide.shapes.add_textbox(Inches(7.2), Inches(6.25), Inches(5.4), Inches(0.8))
    tf = arch_text.text_frame
    p = tf.paragraphs[0]
    p.text = "• Server-side indicator calculations\n• Web Workers for background processing\n• Canvas-based rendering (no TradingView dependency)"
    p.font.size = Pt(11)
    p.font.color.rgb = SUBTITLE_TEXT
    
    return slide

def add_roadmap_slide(prs):
    """Slide 5: Current State and Future Roadmap"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Current State & Future Roadmap"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.LEFT
    
    # Current State - left side
    current_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.5))
    tf = current_title.text_frame
    p = tf.paragraphs[0]
    p.text = "✓ Current State (April 2026)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    current_items = [
        "Layer 4 execution path running with schema-aligned SQL writes",
        "Layer 4 logging with rotating log files",
        "Layer 1 regime pipeline clustering fixed",
        "Layer 5 read-oriented API (10+ endpoints)",
        "Custom HTML5 Canvas chart system deployed",
        "NLP/FinBERT macro intelligence pipeline active",
        "Real-time OANDA WebSocket streaming",
        "30+ technical indicators library",
        "Trade marker visualization (entry/SL/TP)",
    ]
    
    y_pos = 1.75
    for item in current_items:
        # Checkmark
        check = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(0.4), Inches(0.35))
        tf = check.text_frame
        p = tf.paragraphs[0]
        p.text = "✓"
        p.font.size = Pt(12)
        p.font.color.rgb = ACCENT_GREEN
        
        # Text
        item_box = slide.shapes.add_textbox(Inches(0.9), Inches(y_pos), Inches(5.6), Inches(0.35))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_TEXT
        
        y_pos += 0.38
    
    # Future Roadmap - right side
    future_title = slide.shapes.add_textbox(Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.5))
    tf = future_title.text_frame
    p = tf.paragraphs[0]
    p.text = "⟳ Upcoming Improvements"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    roadmap_items = [
        ("Phase 1", "Macro events integration into Layer 3/4", "Short-term"),
        ("Phase 2", "Schema health checks before Layer 4 starts", "Short-term"),
        ("Phase 3", "Documentation sync automation", "Medium-term"),
        ("Phase 4", "Macro sentiment dashboard cards", "Medium-term"),
        ("Phase 5", "Additional chart types (Heikin-Ashi, Renko)", "Long-term"),
        ("Phase 6", "ML pattern detection on charts", "Long-term"),
        ("Phase 7", "Chart templates and saved layouts", "Long-term"),
        ("Phase 8", "Export to SVG/PDF formats", "Long-term"),
    ]
    
    y_pos = 1.75
    for phase, desc, timeline in roadmap_items:
        # Phase box
        phase_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(y_pos), Inches(5.8), Inches(0.55))
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = RGBColor(0x2a, 0x2a, 0x4e)
        phase_box.line.color.rgb = ACCENT_BLUE
        
        # Phase name
        phase_text = slide.shapes.add_textbox(Inches(7.15), Inches(y_pos + 0.05), Inches(1.2), Inches(0.25))
        tf = phase_text.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
        # Description
        desc_text = slide.shapes.add_textbox(Inches(8.4), Inches(y_pos + 0.05), Inches(3.0), Inches(0.25))
        tf = desc_text.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = LIGHT_TEXT
        
        # Timeline
        time_text = slide.shapes.add_textbox(Inches(11.5), Inches(y_pos + 0.05), Inches(1.2), Inches(0.25))
        tf = time_text.text_frame
        p = tf.paragraphs[0]
        p.text = timeline
        p.font.size = Pt(9)
        p.font.color.rgb = SUBTITLE_TEXT
        p.alignment = PP_ALIGN.RIGHT
        
        y_pos += 0.65
    
    # Success metrics box at bottom
    metrics_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.2), Inches(12.333), Inches(1.0))
    metrics_box.fill.solid()
    metrics_box.fill.fore_color.rgb = RGBColor(0x2a, 0x2a, 0x4e)
    metrics_box.line.color.rgb = ACCENT_GREEN
    
    metrics_title = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(12), Inches(0.4))
    tf = metrics_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Success Metrics: Performance < 2s load | 60 FPS | 99.9% uptime | 80%+ test coverage"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    metrics_sub = slide.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(12), Inches(0.4))
    tf = metrics_sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Institutional-grade quality with professional trading platform experience"
    p.font.size = Pt(12)
    p.font.color.rgb = SUBTITLE_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# Generate all slides
add_title_slide(prs)
add_architecture_slide(prs)
add_dataflow_slide(prs)
add_charts_slide(prs)
add_roadmap_slide(prs)

# Save presentation
output_path = "/home/emmanuel/Documents/Scalable_Brain/Scalable_Brain_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
